#!/usr/bin/env python3
"""Generate OpenClaw Deep Dive PPTX with matplotlib-rendered charts."""

import os
import tempfile
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.ticker as mticker
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN_L = Inches(0.8)
MARGIN_T = Inches(0.6)
CONTENT_W = Inches(11.7)

BG      = RGBColor(0x0A, 0x0A, 0x0A)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GRAY    = RGBColor(0xA0, 0xA0, 0xA0)
DIM     = RGBColor(0x72, 0x72, 0x72)
BORDER  = RGBColor(0x30, 0x30, 0x30)
BLUE    = RGBColor(0x4E, 0xA8, 0xDE)
GREEN   = RGBColor(0x1E, 0xD7, 0x60)
RED     = RGBColor(0xEF, 0x44, 0x44)
ORANGE  = RGBColor(0xF5, 0x9E, 0x0B)
PURPLE  = RGBColor(0xA8, 0x55, 0xF7)

CHART_TMP = tempfile.mkdtemp(prefix="oc_pptx_")

MPL = {
    'blue': '#4EA8DE', 'green': '#1ED760', 'red': '#EF4444',
    'orange': '#F59E0B', 'purple': '#A855F7', 'gray': '#727272',
    'dim': '#303030', 'bg': '#0A0A0A', 'light': '#B3B3B3',
}

def save_chart(fig, name, dpi=250):
    path = os.path.join(CHART_TMP, f"{name}.png")
    fig.savefig(path, dpi=dpi, bbox_inches='tight', facecolor=MPL['bg'],
                edgecolor='none', pad_inches=0.2)
    plt.close(fig)
    return path

def add_bg(slide):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG

def add_title(slide, title, subtitle, y=MARGIN_T):
    tf = slide.shapes.add_textbox(MARGIN_L, y, CONTENT_W, Inches(0.6)).text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    if subtitle:
        tf2 = slide.shapes.add_textbox(MARGIN_L, y + Inches(0.55), CONTENT_W, Inches(0.4)).text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(16)
        p2.font.color.rgb = GRAY

def add_slide_num(slide, num):
    tf = slide.shapes.add_textbox(SLIDE_W - Inches(1), SLIDE_H - Inches(0.5), Inches(0.6), Inches(0.3)).text_frame
    p = tf.paragraphs[0]
    p.text = str(num)
    p.font.size = Pt(10)
    p.font.color.rgb = DIM
    p.alignment = PP_ALIGN.RIGHT

def add_text_box(slide, text, x, y, w, h, font_size=14, color=GRAY, bold=False, align=PP_ALIGN.LEFT):
    tf = slide.shapes.add_textbox(x, y, w, h).text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return tf

def add_card(slide, x, y, w, h, title_text, body_text, accent=BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x1A)
    shape.line.fill.solid()
    shape.line.fill.fore_color.rgb = BORDER
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(16)
    tf.margin_right = Pt(16)
    tf.margin_top = Pt(14)
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = accent
    p2 = tf.add_paragraph()
    p2.text = body_text
    p2.font.size = Pt(11)
    p2.font.color.rgb = GRAY
    p2.space_before = Pt(8)

# --- Slide builders ---

def build_slide_01_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_slide_num(slide, 1)
    add_text_box(slide, "220,651", MARGIN_L, Inches(1.5), CONTENT_W, Inches(1.2),
                 font_size=72, color=BLUE, bold=True)
    add_text_box(slide, "GitHub stars in 84 days", MARGIN_L, Inches(2.7), CONTENT_W, Inches(0.4),
                 font_size=16, color=DIM)
    add_title(slide, "OpenClaw Deep Dive", "How the fastest-growing open-source project in history actually works", y=Inches(3.5))
    add_text_box(slide, "A case study built with the PM AI Partner Framework",
                 MARGIN_L, Inches(4.6), CONTENT_W, Inches(0.4), font_size=14, color=DIM)
    add_text_box(slide, "Ahmed Khaled Mohamed  ·  February 2026",
                 MARGIN_L, SLIDE_H - Inches(0.8), CONTENT_W, Inches(0.3), font_size=11, color=DIM)

def build_slide_02_what(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_slide_num(slide, 2)
    add_title(slide, "What Is OpenClaw?", "A self-hosted AI agent that runs locally and acts through your messaging apps")
    desc = (
        "Unlike ChatGPT, OpenClaw executes actions: books flights, manages email, "
        "runs terminal commands, browses the web. Connects through WhatsApp, Telegram, "
        "Slack, Discord, Signal, iMessage. Data stays local. No subscription. MIT licensed.\n\n"
        "Created by Peter Steinberger (PSPDFKit founder). Launched November 2025."
    )
    add_text_box(slide, desc, MARGIN_L, Inches(1.6), Inches(5.5), Inches(3), font_size=14, color=GRAY)
    cards = [
        ("6 messaging platforms", "WhatsApp, Telegram, Slack, Discord, Signal, iMessage"),
        ("100% local", "Runs on your machine, data never leaves"),
        ("3,000+ skills", "ClawHub marketplace, community-built"),
        ("Autonomous", "Heartbeat acts while you're away"),
    ]
    for i, (t, b) in enumerate(cards):
        col = i % 2
        row = i // 2
        add_card(slide, Inches(7) + col * Inches(3.2), Inches(1.6) + row * Inches(1.8),
                 Inches(3), Inches(1.5), t, b, GREEN)

def build_slide_03_arch(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_slide_num(slide, 3)
    add_title(slide, "Architecture", "Four components, one control plane")
    fig, ax = plt.subplots(figsize=(12, 5), facecolor=MPL['bg'])
    ax.set_xlim(0, 12)
    ax.set_ylim(0, 6)
    ax.axis('off')
    channels = ['WhatsApp', 'Telegram', 'Slack', 'Discord', 'Signal', 'iMessage']
    for i, ch in enumerate(channels):
        x = 1 + i * 1.7
        ax.add_patch(plt.Rectangle((x, 4.8), 1.4, 0.7, fill=True, facecolor='#1a1a1a',
                                    edgecolor=MPL['blue'], linewidth=1.5, zorder=2))
        ax.text(x + 0.7, 5.15, ch, ha='center', va='center', color=MPL['blue'],
                fontsize=9, fontweight='bold', zorder=3)
    ax.annotate('', xy=(5.5, 4.2), xytext=(5.5, 4.7),
                arrowprops=dict(arrowstyle='->', color=MPL['dim'], lw=2))
    ax.add_patch(plt.Rectangle((2.5, 3.2), 6, 0.9, fill=True, facecolor='#1a1a1a',
                                edgecolor=MPL['green'], linewidth=2, zorder=2))
    ax.text(5.5, 3.65, 'GATEWAY  (Node.js WebSocket hub)', ha='center', va='center',
            color=MPL['green'], fontsize=13, fontweight='bold', zorder=3)
    ax.annotate('', xy=(5.5, 2.6), xytext=(5.5, 3.1),
                arrowprops=dict(arrowstyle='->', color=MPL['dim'], lw=2))
    components = [
        (1.5, 'BRAIN', 'LLM Reasoning', MPL['purple']),
        (4.5, 'HANDS', 'Shell, Files, Browser', MPL['orange']),
        (7.5, 'MEMORY', 'Local Markdown', MPL['blue']),
    ]
    for x, label, desc, color in components:
        ax.add_patch(plt.Rectangle((x, 1.5), 2.5, 0.9, fill=True, facecolor='#1a1a1a',
                                    edgecolor=color, linewidth=1.5, zorder=2))
        ax.text(x + 1.25, 2.05, label, ha='center', va='center', color=color,
                fontsize=11, fontweight='bold', zorder=3)
        ax.text(x + 1.25, 1.75, desc, ha='center', va='center', color=MPL['light'],
                fontsize=8, zorder=3)
    ax.annotate('', xy=(5.5, 0.7), xytext=(5.5, 1.4),
                arrowprops=dict(arrowstyle='->', color=MPL['dim'], lw=2))
    ax.add_patch(plt.Rectangle((3.5, 0.1), 4, 0.7, fill=True, facecolor='#1a1a1a',
                                edgecolor=MPL['red'], linewidth=1.5, zorder=2))
    ax.text(5.5, 0.45, 'HEARTBEAT  (Every 30 min, autonomous)', ha='center', va='center',
            color=MPL['red'], fontsize=10, fontweight='bold', zorder=3)
    img = save_chart(fig, 'architecture')
    slide.shapes.add_picture(img, Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.5))

def build_slide_04_stars(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_slide_num(slide, 4)
    add_title(slide, "Star Growth: 0 to 220K in 84 Days", "The fastest-growing software repository in GitHub history")
    dates = ['Nov 24', 'Dec 15', 'Jan 1', 'Jan 15', 'Jan 24', 'Jan 26', 'Jan 27', 'Jan 30', 'Feb 8', 'Feb 17', 'Feb 23']
    stars = [0, 200, 500, 800, 1000, 26000, 40000, 60000, 145000, 200000, 220651]
    fig, ax = plt.subplots(figsize=(11, 4.5), facecolor=MPL['bg'])
    ax.set_facecolor(MPL['bg'])
    ax.fill_between(range(len(dates)), stars, alpha=0.15, color=MPL['blue'])
    ax.plot(range(len(dates)), stars, color=MPL['blue'], linewidth=3, marker='o',
            markersize=8, markerfacecolor=MPL['blue'], markeredgecolor=MPL['blue'])
    ax.set_xticks(range(len(dates)))
    ax.set_xticklabels(dates, fontsize=10, color=MPL['light'])
    ax.yaxis.set_major_formatter(mticker.FuncFormatter(lambda v, _: f'{int(v/1000)}K' if v >= 1000 else str(int(v))))
    ax.tick_params(colors=MPL['light'], labelsize=10)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['bottom'].set_color(MPL['dim'])
    ax.spines['left'].set_color(MPL['dim'])
    ax.grid(axis='y', color=MPL['dim'], linewidth=0.5, alpha=0.5)
    ax.annotate('+25,310\nin one day', xy=(5, 26000), xytext=(6.5, 80000),
                color=MPL['orange'], fontsize=10, fontweight='bold',
                arrowprops=dict(arrowstyle='->', color=MPL['orange'], lw=1.5))
    ax.annotate('200K\n(84 days)', xy=(9, 200000), xytext=(7.5, 180000),
                color=MPL['green'], fontsize=10, fontweight='bold',
                arrowprops=dict(arrowstyle='->', color=MPL['green'], lw=1.5))
    img = save_chart(fig, 'stars')
    slide.shapes.add_picture(img, Inches(0.5), Inches(1.6), Inches(12.3), Inches(5.2))

def build_slide_05_ecosystem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_slide_num(slide, 5)
    add_title(slide, "Ecosystem & Composition", "370 contributors, 42K forks, 50 releases, 5 language reimplementations")
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(12, 4), facecolor=MPL['bg'])
    langs = ['TypeScript', 'Swift', 'Other']
    pcts = [85.2, 11.0, 3.8]
    colors = [MPL['blue'], MPL['orange'], '#303030']
    wedges, texts, autotexts = ax1.pie(pcts, labels=langs, colors=colors, autopct='%1.1f%%',
                                        startangle=90, textprops={'color': MPL['light'], 'fontsize': 11})
    for at in autotexts:
        at.set_color('#ffffff')
        at.set_fontsize(10)
    ax1.set_title('Language Composition', color=MPL['light'], fontsize=13, pad=10)
    projects = ['OpenClaw\n(TS)', 'ZeroClaw\n(Rust)', 'PicoClaw\n(Go)', 'NanoClaw\n(Py)', 'TinyClaw\n(Shell)']
    proj_colors = [MPL['blue'], MPL['orange'], MPL['green'], MPL['purple'], MPL['light']]
    bars = ax2.bar(range(5), [220651, 12000, 8500, 6200, 3100], color=proj_colors, width=0.6)
    ax2.set_xticks(range(5))
    ax2.set_xticklabels(projects, fontsize=9, color=MPL['light'])
    ax2.set_facecolor(MPL['bg'])
    ax2.spines['top'].set_visible(False)
    ax2.spines['right'].set_visible(False)
    ax2.spines['bottom'].set_color(MPL['dim'])
    ax2.spines['left'].set_color(MPL['dim'])
    ax2.tick_params(colors=MPL['light'], labelsize=9)
    ax2.yaxis.set_major_formatter(mticker.FuncFormatter(lambda v, _: f'{int(v/1000)}K'))
    ax2.set_title('Ecosystem Stars', color=MPL['light'], fontsize=13, pad=10)
    ax2.grid(axis='y', color=MPL['dim'], linewidth=0.5, alpha=0.5)
    img = save_chart(fig, 'ecosystem')
    slide.shapes.add_picture(img, Inches(0.5), Inches(1.6), Inches(12.3), Inches(5))

def build_slide_06_rebrand(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_slide_num(slide, 6)
    add_title(slide, "The Wildest Week in Open Source", "Three names, one crypto scam, $16M in fake tokens")
    events = [
        ("Nov 2025", "Clawdbot launches", "A play on Anthropic's 'Claude'", BLUE),
        ("Jan 27, 2026", "Rebrands to Moltbot", "Anthropic sends trademark request", ORANGE),
        ("Jan 28, 2026", "Accounts hijacked in ~10 seconds", "Scammers launch $CLAWD token on Solana", RED),
        ("Hours later", "$16M market cap, then -90% crash", "Steinberger denies involvement, token collapses", RED),
        ("Jan 29, 2026", "Final rebrand to OpenClaw", "No abandoned accounts this time", GREEN),
    ]
    for i, (date, event, detail, color) in enumerate(events):
        y = Inches(1.8) + i * Inches(1.0)
        add_text_box(slide, date, MARGIN_L, y, Inches(2), Inches(0.3), font_size=12, color=DIM)
        add_text_box(slide, event, MARGIN_L + Inches(2.2), y, Inches(5), Inches(0.3),
                     font_size=14, color=color, bold=True)
        add_text_box(slide, detail, MARGIN_L + Inches(2.2), y + Inches(0.3), Inches(5), Inches(0.3),
                     font_size=11, color=GRAY)
    add_card(slide, Inches(8), Inches(1.8), Inches(4.5), Inches(4.5),
             "The Lesson",
             "In 2026, a naming decision on an open-source project can trigger "
             "a $16M financial event in under 24 hours. The intersection of open source, "
             "crypto speculation, and social media impersonation created a new risk category.",
             RED)

def build_slide_07_security(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_slide_num(slide, 7)
    add_title(slide, "The Security Wake-Up Call", "6 CVEs, 900 malicious skills, 40,000 exposed instances")
    cves = [
        ("CVE-2026-25593", "Critical", "Unauthenticated command execution"),
        ("CVE-2026-25253", "High 8.8", "One-click RCE via malicious URLs"),
        ("CVE-2026-26323", "High", "CI/supply-chain injection"),
        ("CVE-2026-26327", "Medium", "Gateway impersonation (mDNS)"),
        ("CVE-2026-26317", "Medium", "Privilege escalation"),
        ("CVE-2026-26329", "Medium", "Token replay attacks"),
    ]
    for i, (cve, sev, desc) in enumerate(cves):
        y = Inches(1.7) + i * Inches(0.65)
        sev_color = RED if 'Critical' in sev else ORANGE if 'High' in sev else BLUE
        add_text_box(slide, cve, MARGIN_L, y, Inches(2.2), Inches(0.3), font_size=11, color=DIM)
        add_text_box(slide, sev, MARGIN_L + Inches(2.3), y, Inches(1.2), Inches(0.3),
                     font_size=11, color=sev_color, bold=True)
        add_text_box(slide, desc, MARGIN_L + Inches(3.8), y, Inches(3.5), Inches(0.3),
                     font_size=12, color=GRAY)
    add_card(slide, Inches(8.2), Inches(1.7), Inches(4.3), Inches(2.5),
             "The Lethal Trifecta",
             "1. Access to private data (files, emails, messages)\n"
             "2. Exposure to untrusted content (web browsing)\n"
             "3. External communication (send messages, API calls)\n\n"
             "A compromised agent can exfiltrate data through channels you'd never check.",
             RED)
    add_card(slide, Inches(8.2), Inches(4.6), Inches(2), Inches(1.5),
             "~900", "Malicious skills\n(20% of ClawHub)", RED)
    add_card(slide, Inches(10.5), Inches(4.6), Inches(2), Inches(1.5),
             "40,000+", "Exposed instances\n(35% vulnerable)", ORANGE)

def build_slide_08_insights(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_slide_num(slide, 8)
    add_title(slide, "Why OpenClaw Exploded: 3 Product Insights",
              "The innovation isn't the AI model — it's the interaction model")
    insights = [
        ("01", "Local-First Is the Killer Feature",
         "Runs on your machine, data never leaves. Users bought Mac Minis as dedicated agent machines — Apple stores sold out.",
         GREEN),
        ("02", "Messaging-as-Interface",
         "Meets you in WhatsApp, Telegram, Slack. No new UI. Distribution through existing channels beats purpose-built interfaces.",
         BLUE),
        ("03", "The Heartbeat Creates Attachment",
         "Every 30 min, checks tasks and acts autonomously. Transforms a tool into an assistant. Chatbots respond. Agents initiate.",
         PURPLE),
    ]
    for i, (num, title, body, color) in enumerate(insights):
        x = MARGIN_L + i * Inches(4)
        add_text_box(slide, num, x, Inches(1.8), Inches(1), Inches(0.6),
                     font_size=36, color=BORDER, bold=True)
        add_card(slide, x, Inches(2.5), Inches(3.6), Inches(3.5), title, body, color)

def build_slide_09_framework(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_slide_num(slide, 9)
    add_title(slide, "How This Was Built: The Framework in Action",
              "4 skills, ~3 hours, from zero knowledge to complete analysis")
    skills = [
        ("/technical-analyst", "Map the Architecture",
         "Trace message flow, identify components, diagram the system", BLUE),
        ("/data-analyst", "Collect the Numbers",
         "Star growth, contributors, language breakdown, ecosystem stats", GREEN),
        ("/devil-advocate", "Challenge the Narrative",
         "CVE analysis, permission model risks, malicious skills data", RED),
        ("/builder", "Build the Deliverables",
         "HTML presentation, PPTX deck, Substack article", PURPLE),
    ]
    for i, (cmd, title, desc, color) in enumerate(skills):
        y = Inches(1.7) + i * Inches(1.25)
        add_text_box(slide, cmd, MARGIN_L, y, Inches(2.2), Inches(0.3),
                     font_size=13, color=color, bold=True)
        add_text_box(slide, title, MARGIN_L + Inches(2.5), y, Inches(3), Inches(0.3),
                     font_size=14, color=WHITE, bold=True)
        add_text_box(slide, desc, MARGIN_L + Inches(2.5), y + Inches(0.35), Inches(4.5), Inches(0.5),
                     font_size=11, color=GRAY)
    add_card(slide, Inches(8), Inches(1.7), Inches(4.5), Inches(4.5),
             "The Compounding Effect",
             "Technical analysis told us what to measure.\n"
             "Data collection revealed growth inflection points.\n"
             "Devil's advocate found the security counter-narrative.\n"
             "Builder mode turned it all into shareable formats.\n\n"
             "Structured thinking compounds. Each mode contributes something the others can't.",
             GREEN)

def build_slide_10_takeaways(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_slide_num(slide, 10)
    add_title(slide, "Key Takeaways", None)
    takeaways = [
        ("01", "You don't need to be a developer to understand complex technical projects.", BLUE),
        ("02", "OpenClaw's innovation is the distribution, not the AI.", GREEN),
        ("03", "Speed and security are inversely correlated at scale.", RED),
        ("04", "Open-source AI agents have new risk categories.", ORANGE),
        ("05", "AI works best as a thinking partner when you give it structure.", PURPLE),
    ]
    for i, (num, text, color) in enumerate(takeaways):
        y = Inches(1.5) + i * Inches(1.0)
        add_text_box(slide, num, MARGIN_L, y, Inches(0.6), Inches(0.4),
                     font_size=18, color=color, bold=True)
        add_text_box(slide, text, MARGIN_L + Inches(0.8), y, Inches(10), Inches(0.6),
                     font_size=15, color=WHITE, bold=True)
    add_text_box(slide, "Framework: github.com/ahmedkhaledmohamed/PM-AI-Partner-Framework",
                 MARGIN_L, SLIDE_H - Inches(1), CONTENT_W, Inches(0.3), font_size=11, color=DIM)


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    build_slide_01_title(prs)
    build_slide_02_what(prs)
    build_slide_03_arch(prs)
    build_slide_04_stars(prs)
    build_slide_05_ecosystem(prs)
    build_slide_06_rebrand(prs)
    build_slide_07_security(prs)
    build_slide_08_insights(prs)
    build_slide_09_framework(prs)
    build_slide_10_takeaways(prs)

    out = os.path.join(os.path.dirname(__file__), 'openclaw-deep-dive.pptx')
    prs.save(out)
    print(f"Saved: {out}")
    print(f"Chart images: {CHART_TMP}")


if __name__ == '__main__':
    main()
